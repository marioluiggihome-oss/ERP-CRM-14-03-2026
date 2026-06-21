#!/usr/bin/env python3
"""Genera el PowerPoint comercial de Luiggi Home ERP (deck de venta SaaS).

Uso: python3 scripts/generate_pitch_deck.py
Salida: docs/Luiggi_Home_ERP.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Paleta de marca ──────────────────────────────────────────────────────────
ORANGE = RGBColor(0xEA, 0x58, 0x0C)   # naranja Luiggi
INDIGO = RGBColor(0x31, 0x2E, 0x81)   # indigo oscuro
DARK   = RGBColor(0x0F, 0x17, 0x2A)   # casi negro (slate-900)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x47, 0x55, 0x69)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def _set(p, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Arial"


def _bar(slide):
    """Barra naranja decorativa arriba."""
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.18))
    sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE
    sh.line.fill.background()


def content_slide(title, subtitle, bullets):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    _bar(s)
    t = _box(s, 0.7, 0.5, 12.0, 1.1)
    _set(t.text_frame.paragraphs[0], title, 30, INDIGO, bold=True)
    if subtitle:
        p = t.text_frame.add_paragraph()
        _set(p, subtitle, 15, ORANGE, bold=True)
    body = _box(s, 0.8, 1.9, 11.7, 5.0)
    tf = body.text_frame
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set(p, "•  " + b, 19, SLATE)
        p.space_after = Pt(10)
    return s


# 1 ── Portada
s = prs.slides.add_slide(BLANK); _bg(s, DARK); _bar(s)
t = _box(s, 0.9, 2.4, 11.5, 2.6)
_set(t.text_frame.paragraphs[0], "LUIGGI HOME", 54, WHITE, bold=True)
p = t.text_frame.add_paragraph()
_set(p, "El ERP + CRM hecho para fabricantes de muebles y cocinas", 24, ORANGE, bold=True)
p = t.text_frame.add_paragraph()
_set(p, "Presupuesto por tarifa real · Despiece · Fábrica · Finanzas · IA", 16, RGBColor(0xCB,0xD5,0xE1))
f = _box(s, 0.9, 6.6, 11.5, 0.5)
_set(f.text_frame.paragraphs[0], "Propuesta comercial — Modelo SaaS", 12, RGBColor(0x94,0xA3,0xB8))

# 2 ── El problema
content_slide("El problema", "Lo que hoy cuesta tiempo y dinero", [
    "Presupuestar a mano: lento y con errores de tarifa.",
    "Despieces y optimización de tableros calculados a ojo → desperdicio de material.",
    "Fábrica y red comercial sin control centralizado.",
    "Herramientas genéricas que no entienden el flujo mueble/cocina.",
])

# 3 ── La solución
content_slide("La solución: Luiggi Home", "Un solo sistema, de principio a fin", [
    "CRM → Presupuesto por tarifa real → Pedido → Fábrica → Finanzas.",
    "Catálogos reales (ZC y MV) con tarifa por zonas T1–T21.",
    "Despiece y optimización de tableros automáticos.",
    "IA propia: render 3D, cocinas 3D y digitalización de planos.",
])

# 4 ── Módulos
content_slide("Módulos clave", "Todo integrado", [
    "Presupuestador de cocinas (12 familias) y configurador de armarios a medida.",
    "Despiece industrial: lista de corte, cantos, herrajes y lista de compra.",
    "Portal de Fábrica con órdenes de fabricación automáticas.",
    "Pedidos, facturación, rentabilidad, inventario y montajes.",
])

# 5 ── Diferenciador IA
content_slide("Diferenciador: IA propia", "Lo que no tiene un ERP genérico", [
    "Render fotorrealista del mueble/cocina fiel a la configuración.",
    "Diseño de cocinas 3D con IA.",
    "Digitalizador de planos: de una foto a un presupuesto en segundos.",
    "Dictado por voz para configurar más rápido.",
])

# 6 ── Tarifa real + despiece
content_slide("Tarifa real + despiece", "Ahorro directo de tiempo y material", [
    "Precio exacto por zona de tarifa (T1–T21), sin hojas de cálculo.",
    "Optimización de tableros con cálculo de merma.",
    "Anulación de líneas (p. ej. puertas que se piden fuera).",
    "Menos errores = menos retrabajos en fábrica.",
])

# 7 ── Planes y precios
s = content_slide("Planes y precios", "Escalado por tipo de empresa", [])
from pptx.util import Inches as IN
rows = [
    ("Plan", "Empresa", "Usuarios", "€/mes"),
    ("Starter", "Autónomo / tienda", "1–3", "99–149"),
    ("Pro", "PYME de cocinas", "hasta 10", "299–399"),
    ("Business", "Fabricante + comercial", "hasta 25", "699–899"),
    ("Premium IA", "Fabricante / grupo", "25+", "1.290–1.890"),
    ("Enterprise", "Grupo / franquicia", "50+", "desde 2.500"),
]
tbl = s.shapes.add_table(len(rows), 4, IN(0.8), IN(2.1), IN(11.7), IN(4.6)).table
for ci, w in enumerate((2.6, 4.3, 2.3, 2.5)):
    tbl.columns[ci].width = IN(w)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        r = para.runs[0]; r.font.size = Pt(14); r.font.name = "Arial"
        if ri == 0:
            r.font.bold = True; r.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
        else:
            r.font.color.rgb = SLATE
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            if ci == 0:
                r.font.bold = True; r.font.color.rgb = ORANGE

# 8 ── ROI
content_slide("ROI: se paga solo", "Venta anclada al ahorro", [
    "Cada presupuesto y cada despiece ahorra horas de trabajo.",
    "Menos desperdicio de tablero por optimización real.",
    "Add-ons: marca blanca, créditos IA, onboarding, SLA premium.",
    "Facturación anual con 2 meses gratis (~17 % de descuento).",
])

# 9 ── Marca blanca / canal
content_slide("Marca blanca y canal", "Crecimiento por distribuidores", [
    "El distribuidor de tablero/herrajes lo revende con su logo y dominio.",
    "Onboarding + migración de datos como servicio.",
    "Nicho España (carpinterías/cocinas) → LATAM.",
    "Demo guiada + prueba de 14 días.",
])

# 10 ── Oportunidad
content_slide("Oportunidad de mercado", "Proyección ilustrativa", [
    "40 empresas, mix de planes → ARPU ~450 €/empresa·mes.",
    "ARR ≈ 216.000 €.",
    "A múltiplos SaaS (3–6× ARR): valoración ~650.000 – 1,3 M €.",
    "Churn objetivo <3 %/mes · LTV/CAC >3×.",
])

# 11 ── Roadmap
content_slide("Roadmap", "Hacia dónde va", [
    "IA aún más potente para render y cocinas 3D.",
    "Multi-empresa y marca blanca completa.",
    "App móvil para comercial y montadores.",
    "Más catálogos y tarifas integradas.",
])

# 12 ── Seguridad y soporte
content_slide("Seguridad y soporte", "Listo para empresa", [
    "Permisos por capacidad y por rol (lo que no está activo, no se ve).",
    "Candado de costes: cada usuario ve solo lo que debe.",
    "Backups, mantenimiento y telemetría.",
    "Soporte y SLA según plan.",
])

# 13 ── Cierre / CTA
s = prs.slides.add_slide(BLANK); _bg(s, INDIGO); _bar(s)
t = _box(s, 0.9, 2.6, 11.5, 2.4)
_set(t.text_frame.paragraphs[0], "Empieza hoy", 44, WHITE, bold=True)
p = t.text_frame.add_paragraph()
_set(p, "Prueba 14 días · Demo guiada · Migración incluida", 22, ORANGE, bold=True)
p = t.text_frame.add_paragraph()
_set(p, "Luiggi Home — software@luiggihome.es", 16, RGBColor(0xCB,0xD5,0xE1))

import os
os.makedirs("docs", exist_ok=True)
out = "docs/Luiggi_Home_ERP.pptx"
prs.save(out)
print("OK", out, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
